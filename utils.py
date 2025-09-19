# utils.py
import os
import io
import tempfile
from supabase import create_client, Client
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

# File readers
import docx2txt
import PyPDF2
import openpyxl
from pptx import Presentation

# -------------------------
# Shared Supabase constants
# -------------------------
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_KEY")
SUPABASE_SERVICE_KEY = os.getenv("SUPABASE_SERVICE_KEY")
DOCS_BUCKET = os.getenv("DOCS_BUCKET", "chatbot-docs")  # Default fallback
SUMMARY_BUCKET = os.getenv("SUMMARY_BUCKET", "summarized-text")  # Default fallback
SUMMARY_FILENAME = os.getenv("SUMMARY_FILENAME", "summarized_text.md")  # Default fallback

# -------------------------
# Create Supabase client
# -------------------------
def get_supabase_client() -> Client:
    return create_client(SUPABASE_URL, SUPABASE_KEY)

def get_supabase_service_client() -> Client:
    """Get Supabase client with service role for admin operations"""
    return create_client(SUPABASE_URL, SUPABASE_SERVICE_KEY)

# -------------------------
# Fetch summarized text
# -------------------------
async def fetch_summarized_text() -> str:
    """
    Downloads the summarized_text.md file from Supabase.
    Returns the content as a string, or empty string if missing.
    """
    supabase = get_supabase_client()
    try:
        resp = supabase.storage.from_(SUMMARY_BUCKET).download(SUMMARY_FILENAME)
        return resp.decode("utf-8")
    except Exception as e:
        print(f"[utils] No summarized_text.md found: {e}")
        return ""

# -------------------------
# Extract text helpers
# -------------------------
def extract_text(filename: str, file_bytes: bytes) -> str:
    text = ""
    if filename.endswith(".pdf"):
        reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    elif filename.endswith(".docx"):
        # Use a Windows-compatible temp path
        import tempfile
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as temp_file:
            temp_file.write(file_bytes)
            temp_path = temp_file.name
        
        try:
            text = docx2txt.process(temp_path)
        finally:
            # Clean up the temp file
            if os.path.exists(temp_path):
                os.remove(temp_path)
    elif filename.endswith(".pptx"):
        prs = Presentation(io.BytesIO(file_bytes))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif filename.endswith(".xlsx"):
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes))
        for sheet in wb:
            for row in sheet.iter_rows(values_only=True):
                text += " ".join([str(cell) if cell else "" for cell in row]) + "\n"
    elif filename.endswith(".txt") or filename.endswith(".md"):
        text = file_bytes.decode("utf-8", errors="ignore")
    else:
        text = f"[utils] Unsupported file type: {filename}\n"
    return text.strip()

def discover_files_by_common_patterns():
    """
    Discover files by trying common filename patterns when bucket listing fails.
    This is a fallback method to find files in the chatbot-docs bucket.
    """
    discovered_files = []
    supabase_service = get_supabase_service_client()
    
    # Common file extensions we support
    extensions = ['docx', 'pdf', 'txt', 'xlsx', 'pptx', 'doc', 'xls', 'ppt']
    
    # Try some common timestamp ranges (files are usually named with timestamps)
    import time
    current_time = int(time.time())
    
    # Check recent uploads (last 90 days)
    for days_ago in range(90):
        timestamp = current_time - (days_ago * 24 * 60 * 60)
        
        # Try different hash patterns (the part after underscore)
        # Common patterns: 8-12 character hex strings
        test_hashes = [
            '536ae6173c20',  # Known working hash
            '801b0240f7ce',  # Another known hash
            format(timestamp % 0xffffffffffff, '012x'),  # Generate from timestamp
            format((timestamp * 31) % 0xffffffffffff, '012x'),  # Variation
        ]
        
        for ext in extensions:
            for hash_part in test_hashes:
                test_filename = f"{timestamp}_{hash_part}.{ext}"
                try:
                    test_download = supabase_service.storage.from_(DOCS_BUCKET).download(test_filename)
                    if test_filename not in discovered_files:
                        discovered_files.append(test_filename)
                        print(f"[utils] Discovered file: {test_filename} ({len(test_download)} bytes)")
                except:
                    continue  # File doesn't exist
                    
                # Don't try too many combinations to avoid excessive API calls
                if len(discovered_files) >= 10:
                    break
            if len(discovered_files) >= 10:
                break
        if len(discovered_files) >= 10:
            break
    
    return discovered_files

# -------------------------
# Compile docs and store (overwrite safe)
# -------------------------
async def summarize_and_store():
    """
    Reads documents from Supabase DOCS_BUCKET,
    extracts text without AI summarization,
    and saves compiled summarized_text.md into SUMMARY_BUCKET.
    """
    supabase = get_supabase_client()

    # Try to list all files in chatbot-docs using both regular and service clients
    files = []
    listing_failed = False
    
    # First try with regular client
    try:
        files = supabase.storage.from_(DOCS_BUCKET).list()
        print(f"[utils] Found {len(files)} files via regular client bucket listing")
        
        if len(files) > 0:
            listing_failed = False
        else:
            print(f"[utils] Regular client returned empty - trying service client")
            listing_failed = True
            
    except Exception as e:
        print(f"[utils] Regular client bucket listing failed: {e}")
        listing_failed = True
    
    # If regular client failed, try with service client (better permissions)
    if listing_failed:
        try:
            supabase_service = get_supabase_service_client()
            files = supabase_service.storage.from_(DOCS_BUCKET).list()
            print(f"[utils] Found {len(files)} files via service client bucket listing")
            
            if len(files) > 0:
                listing_failed = False
            else:
                print(f"[utils] Service client also returned empty - using fallback")
                listing_failed = True
                
        except Exception as e:
            print(f"[utils] Service client bucket listing also failed: {e}")
            listing_failed = True
        
    # If both clients failed, use fallback approaches
    if listing_failed:
        print(f"[utils] Using fallback approaches...")
        files = []  # Reset files array
        all_possible_files = set()  # Use set to avoid duplicates
        
        # Method 1: Try known existing files
        known_files = [
            "1758076583_536ae6173c20.docx",
            "1758112830_801b0240f7ce.docx"
        ]
        
        for file_name in known_files:
            all_possible_files.add(file_name)
        
        # Method 2: Try to discover files using common patterns
        try:
            discovered_files = discover_files_by_common_patterns()
            for filename in discovered_files:
                all_possible_files.add(filename)
            print(f"[utils] Pattern discovery found {len(discovered_files)} additional files")
        except Exception as e:
            print(f"[utils] Pattern discovery failed: {e}")
        
        print(f"[utils] Fallback: Testing {len(all_possible_files)} potential files...")
        
        for file_name in all_possible_files:
            try:
                # Try with service client for downloads (better permissions)
                supabase_service = get_supabase_service_client()
                test_download = supabase_service.storage.from_(DOCS_BUCKET).download(file_name)
                files.append({"name": file_name})
                print(f"[utils] Confirmed file: {file_name} ({len(test_download)} bytes)")
            except Exception as file_e:
                print(f"[utils] File not accessible: {file_name} - {file_e}")

    if not files:
        print("[utils] No files found in chatbot-docs bucket after trying all methods")
        return

    compiled_parts = []

    for file in files:
        file_name = file["name"]
        print(f"[utils] Processing {file_name}...")

        try:
            raw_bytes = supabase.storage.from_(DOCS_BUCKET).download(file_name)
            text = extract_text(file_name, raw_bytes)
            if text:
                compiled_parts.append(f"## {file_name}\n\n{text}")
                print(f"[utils] Extracted {len(text)} characters from {file_name}")
        except Exception as e:
            print(f"[utils] Failed to process {file_name}: {e}")

    if not compiled_parts:
        print("[utils] No content extracted from any files")
        return

    # Join all extracted text
    final_summary = "\n\n---\n\n".join(compiled_parts)
    print(f"[utils] Compiled summary: {len(final_summary)} characters total")

    # Delete old summary if exists (use service client for write operations)
    supabase_service = get_supabase_service_client()
    try:
        supabase_service.storage.from_(SUMMARY_BUCKET).remove([SUMMARY_FILENAME])
        print(f"[utils] Removed old {SUMMARY_FILENAME}")
    except Exception as e:
        print(f"[utils] No old summary to delete: {e}")

    # Upload new one (use service client for write operations)
    try:
        supabase_service.storage.from_(SUMMARY_BUCKET).upload(
            SUMMARY_FILENAME,
            final_summary.encode("utf-8"),
            {
                "content-type": "text/markdown"
            }
        )
        print(f"[utils] ✅ Successfully uploaded {SUMMARY_FILENAME} to {SUMMARY_BUCKET}")
        print(f"[utils] Summary file size: {len(final_summary)} characters")
    except Exception as e:
        print(f"[utils] Failed to upload summary: {e}")


