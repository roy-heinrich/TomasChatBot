# summarizer.py
import os
import logging
import asyncio
import httpx
import io
from utils import get_supabase_client, DOCS_BUCKET, SUMMARY_BUCKET, SUMMARY_FILENAME

# File type extraction libraries
import docx2txt
import PyPDF2
import openpyxl
from pptx import Presentation

logger = logging.getLogger("summarizer")
OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")
OPENROUTER_MODEL = "google/gemma-3n-e2b-it:free"

# -------------------------
# File text extraction
# -------------------------
def extract_text_from_file(file_name: str, raw_bytes: bytes) -> str:
    ext = file_name.lower().split(".")[-1]
    text = ""
    try:
        if ext == "txt":
            text = raw_bytes.decode("utf-8", errors="ignore")
        elif ext == "docx":
            f = io.BytesIO(raw_bytes)
            text = docx2txt.process(f)
        elif ext == "pdf":
            f = io.BytesIO(raw_bytes)
            reader = PyPDF2.PdfReader(f)
            text = "\n".join([p.extract_text() or "" for p in reader.pages])
        elif ext in ["xlsx", "xls"]:
            f = io.BytesIO(raw_bytes)
            wb = openpyxl.load_workbook(f, data_only=True)
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    text += " ".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
        elif ext == "pptx":
            f = io.BytesIO(raw_bytes)
            prs = Presentation(f)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        else:
            logger.warning(f"⚠ Unsupported file type {ext} for {file_name}")
    except Exception as e:
        logger.exception(f"❌ Failed to extract text from {file_name}: {e}")
    return text

# -------------------------
# OpenRouter call with retry
# -------------------------
async def call_openrouter_with_retry(client, chunk: str, retries=5, delay=5):
    system_prompt = (
        "Summarize this document into clear Markdown with sections: "
        "Personnel, Departments, Events, Locations, Contact Info. Preserve names and titles."
    )
    for attempt in range(1, retries + 1):
        try:
            resp = await client.post(
                "https://openrouter.ai/api/v1/chat/completions",
                headers={"Authorization": f"Bearer {OPENROUTER_API_KEY}"},
                json={
                    "model": OPENROUTER_MODEL,
                    "messages": [
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": chunk},
                    ],
                },
            )
            if resp.status_code == 200:
                return resp
            elif resp.status_code == 429:
                wait_time = delay * attempt
                logger.warning(f"⚠ Rate limit hit. Retry {attempt}/{retries} after {wait_time}s")
                await asyncio.sleep(wait_time)
            else:
                logger.error(f"❌ OpenRouter failed with status {resp.status_code}")
                break
        except Exception as e:
            logger.exception(f"❌ Exception during OpenRouter request: {e}")
            await asyncio.sleep(delay)
    return None

# -------------------------
# Main summarization
# -------------------------
async def summarize_and_store(
    output_filename: str = SUMMARY_FILENAME,
    chunk_delay: float = 2.0,
    max_consecutive_failures: int = 3
) -> str:
    """
    Summarizes all docs in DOCS_BUCKET and saves result to SUMMARY_BUCKET.
    Includes retries, queueing/delay, and fail-safe stopping if API consistently fails.
    """
    try:
        supabase = get_supabase_client()

        logger.info("⬇ Listing all files in DOCS_BUCKET...")
        files_list = supabase.storage.from_(DOCS_BUCKET).list()
        if not files_list:
            logger.warning("⚠ No files found in DOCS_BUCKET")
            return ""

        # Download and extract text from all files
        combined_text = ""
        for file_obj in files_list:
            file_name = file_obj["name"]
            logger.info(f"⬇ Downloading {file_name}...")
            raw_data = supabase.storage.from_(DOCS_BUCKET).download(file_name)
            if raw_data:
                file_text = extract_text_from_file(file_name, raw_data)
                if file_text.strip():
                    combined_text += file_text + "\n\n"

        if not combined_text.strip():
            logger.warning("⚠ All files were empty after extraction")
            return ""

        # Determine chunk size
        total_length = len(combined_text)
        CHUNK_SIZE = 3000 if total_length < 15000 else 1500
        chunks = [combined_text[i:i+CHUNK_SIZE] for i in range(0, len(combined_text), CHUNK_SIZE)]
        summaries = []

        consecutive_failures = 0
        async with httpx.AsyncClient(timeout=120.0) as client:
            for idx, chunk in enumerate(chunks):
                logger.info(f"📝 Summarizing chunk {idx+1}/{len(chunks)} ({len(chunk)} chars)...")
                resp = await call_openrouter_with_retry(client, chunk)
                if resp is None:
                    logger.error(f"❌ OpenRouter failed on chunk {idx+1}")
                    consecutive_failures += 1
                    if consecutive_failures >= max_consecutive_failures:
                        logger.error("❌ Too many consecutive chunk failures. Aborting summarization.")
                        return ""
                    continue
                consecutive_failures = 0
                data = resp.json()
                summary_part = data["choices"][0]["message"]["content"].strip()
                if summary_part:
                    summaries.append(summary_part)
                else:
                    logger.warning(f"⚠ Chunk {idx+1} returned empty summary")
                # Queue/delay to avoid hitting rate limits
                await asyncio.sleep(chunk_delay)

        final_summary = "\n\n".join(summaries)
        if not final_summary.strip():
            logger.warning("⚠ Final summary is empty after processing all chunks")
            return ""

        # Upload summary to SUMMARY_BUCKET
        logger.info(f"⬆ Uploading {output_filename} to SUMMARY_BUCKET...")
        supabase.storage.from_(SUMMARY_BUCKET).upload(
            output_filename,
            final_summary.encode("utf-8"),
            {"content-type": "text/markdown", "upsert": "true"},
        )

        logger.info("✅ Summarization complete and stored in Supabase.")
        return final_summary

    except Exception as e:
        logger.exception("❌ Error during summarization pipeline")
        return ""
